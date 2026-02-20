"""PPTX generator for AI Transformation Blueprint presentations."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from .models import TransformationBlueprint
from .maturity import MATURITY_LEVELS

# ---------------------------------------------------------------------------
# Colour constants (matching GVRN-AI brand)
# ---------------------------------------------------------------------------
_BG = RGBColor(0x1A, 0x1A, 0x2E)
_GREEN_ACC = RGBColor(0x00, 0xC9, 0x7B)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
_MID_GREY = RGBColor(0x88, 0x88, 0x88)
_DARK_CARD = RGBColor(0x22, 0x22, 0x3A)
_TABLE_ROW1 = RGBColor(0x20, 0x20, 0x38)
_TABLE_ROW2 = RGBColor(0x28, 0x28, 0x44)
_RED_RISK = RGBColor(0xFF, 0x44, 0x44)
_ORANGE = RGBColor(0xFF, 0xA5, 0x00)
_BLUE = RGBColor(0x44, 0x88, 0xFF)

_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers (adapted from audit_toolkit/pptx_gen.py)
# ---------------------------------------------------------------------------

def _set_slide_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _BG

def _add_textbox(slide, left, top, width, height, text,
                 font_size=14, colour=_WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = colour
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def _add_table(slide, rows_data, left, top, width, col_widths=None, row_font_size=10):
    num_rows = len(rows_data)
    num_cols = len(rows_data[0]) if rows_data else 0
    if num_rows == 0 or num_cols == 0:
        return None
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, Inches(0.3 * num_rows))
    table = table_shape.table
    # Remove borders
    for row in table.rows:
        for cell in row.cells:
            for edge in ["top", "bottom", "left", "right"]:
                border = cell._tc.get_or_add_tcPr()
                ln = border.find(qn(f"a:ln{edge.capitalize()}"))
                if ln is None:
                    ln = border.makeelement(qn(f"a:ln{edge.capitalize()}"), {})
                    border.append(ln)
                ln.set("w", "0")
                no_fill = ln.makeelement(qn("a:noFill"), {})
                ln.append(no_fill)
    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            if i < num_cols:
                table.columns[i].width = w
    # Populate cells
    for r_idx, row_data in enumerate(rows_data):
        is_header = r_idx == 0
        bg_color = _GREEN_ACC if is_header else (_TABLE_ROW1 if r_idx % 2 == 1 else _TABLE_ROW2)
        text_color = _BG if is_header else _LIGHT_GREY
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(row_font_size)
                paragraph.font.color.rgb = text_color
                paragraph.font.bold = is_header
            # Cell background
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
    return table_shape

def _slide_title_bar(slide, title, subtitle=""):
    _add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.5),
                 title, font_size=24, colour=_GREEN_ACC, bold=True)
    if subtitle:
        _add_textbox(slide, Inches(0.6), Inches(0.8), Inches(10), Inches(0.3),
                     subtitle, font_size=12, colour=_MID_GREY)

def _add_card(slide, left, top, width, height, title, lines, accent_colour=_GREEN_ACC, bg_colour=_DARK_CARD):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_colour
    shape.line.fill.background()
    # Title
    _add_textbox(slide, left + Inches(0.15), top + Inches(0.1),
                 width - Inches(0.3), Inches(0.3),
                 title, font_size=11, colour=accent_colour, bold=True)
    # Lines
    y = top + Inches(0.4)
    for line in lines:
        _add_textbox(slide, left + Inches(0.15), y,
                     width - Inches(0.3), Inches(0.25),
                     line, font_size=10, colour=_LIGHT_GREY)
        y += Inches(0.22)
    return shape

def _add_metric_box(slide, left, top, width, height, label, value, colour=_GREEN_ACC):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _DARK_CARD
    shape.line.fill.background()
    _add_textbox(slide, left, top + Inches(0.1), width, Inches(0.3),
                 str(value), font_size=28, colour=colour, bold=True,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, left, top + Inches(0.55), width, Inches(0.3),
                 label, font_size=10, colour=_MID_GREY,
                 alignment=PP_ALIGN.CENTER)
    return shape

def _maturity_colour(level: int) -> RGBColor:
    colours = {0: _MID_GREY, 1: _RED_RISK, 2: _ORANGE, 3: RGBColor(0xFF, 0xFF, 0x00), 4: _BLUE, 5: _GREEN_ACC}
    return colours.get(level, _MID_GREY)


# ---------------------------------------------------------------------------
# Main generator
# ---------------------------------------------------------------------------

def generate_blueprint_pptx(blueprint: TransformationBlueprint, output_path: str | Path) -> Path:
    """Generate a full AI Transformation Blueprint PPTX presentation."""
    output_path = Path(output_path)
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank

    # ------------------------------------------------------------------
    # SLIDE 1: Title
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    # Green accent bar top
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), _SLIDE_W, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = _GREEN_ACC; bar.line.fill.background()
    _add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.8),
                 "AI Transformation Blueprint", font_size=40, colour=_WHITE, bold=True)
    _add_textbox(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(0.5),
                 blueprint.organisation_name, font_size=28, colour=_GREEN_ACC)
    _add_textbox(slide, Inches(0.8), Inches(3.6), Inches(11), Inches(0.5),
                 f"Prepared by {blueprint.assessed_by}  |  {blueprint.assessment_date}",
                 font_size=14, colour=_MID_GREY)
    _add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.3),
                 f"Industry: {blueprint.industry.replace('_', ' ').title()}  |  "
                 f"Employees: {blueprint.employee_count}  |  "
                 f"Current Maturity: {blueprint.overall_current_maturity}/5",
                 font_size=12, colour=_LIGHT_GREY)
    # Bottom bar
    bar2 = slide.shapes.add_shape(1, Inches(0), _SLIDE_H - Inches(0.08), _SLIDE_W, Inches(0.08))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = _GREEN_ACC; bar2.line.fill.background()

    # ------------------------------------------------------------------
    # SLIDE 2: Current Maturity Assessment
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _slide_title_bar(slide, "Current AI Maturity Assessment",
                     f"Overall: {blueprint.overall_current_maturity}/5.0")

    _add_metric_box(slide, Inches(9.5), Inches(0.3), Inches(3.5), Inches(0.9),
                    "Current Maturity", f"{blueprint.overall_current_maturity}/5",
                    colour=_maturity_colour(round(blueprint.overall_current_maturity)))

    if blueprint.department_assessments:
        rows = [("Department", "Current Level", "Maturity Name", "Key Pain Points")]
        for dept in blueprint.department_assessments:
            level_name = MATURITY_LEVELS[min(dept.current_maturity, 5)].name if dept.current_maturity <= 5 else "Unknown"
            pain = "; ".join(dept.pain_points[:2]) if dept.pain_points else "None identified"
            rows.append((dept.department, str(dept.current_maturity), level_name, pain[:60]))
        _add_table(slide, rows, Inches(0.4), Inches(1.5), Inches(12.5),
                   col_widths=[Inches(2), Inches(1.5), Inches(2.5), Inches(6.5)],
                   row_font_size=10)

    # ------------------------------------------------------------------
    # SLIDE 3: Target State Vision
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _slide_title_bar(slide, "Target State Vision",
                     f"Target: {blueprint.overall_target_maturity}/5.0  |  "
                     f"Gap: {blueprint.overall_target_maturity - blueprint.overall_current_maturity:.1f}")

    _add_metric_box(slide, Inches(9.5), Inches(0.3), Inches(3.5), Inches(0.9),
                    "Target Maturity", f"{blueprint.overall_target_maturity}/5",
                    colour=_GREEN_ACC)

    if blueprint.department_assessments:
        rows = [("Department", "Current", "Target", "Gap", "Quick Wins")]
        for dept in blueprint.department_assessments:
            qw = "; ".join(dept.quick_wins[:2]) if dept.quick_wins else ""
            rows.append((dept.department, str(dept.current_maturity), str(dept.target_maturity),
                         str(dept.gap), qw[:55]))
        _add_table(slide, rows, Inches(0.4), Inches(1.5), Inches(12.5),
                   col_widths=[Inches(2), Inches(1.2), Inches(1.2), Inches(1), Inches(7.1)],
                   row_font_size=10)

    # ------------------------------------------------------------------
    # SLIDES 4-8: Department Deep Dives (1 per department)
    # ------------------------------------------------------------------
    for dept in blueprint.department_assessments:
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide)
        _slide_title_bar(slide, f"Department: {dept.department}",
                         f"Maturity: {dept.current_maturity} -> {dept.target_maturity} (Gap: {dept.gap})")

        # Pain points card
        pain_lines = [f"- {p[:70]}" for p in dept.pain_points[:5]]
        _add_card(slide, Inches(0.4), Inches(1.3), Inches(6), Inches(2.5),
                  "CURRENT PAIN POINTS", pain_lines, accent_colour=_RED_RISK)

        # Quick wins card
        qw_lines = [f"- {q[:70]}" for q in dept.quick_wins[:4]]
        _add_card(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(2.5),
                  "QUICK WINS", qw_lines, accent_colour=_GREEN_ACC)

        # AI interventions table
        if dept.ai_interventions:
            rows = [("AI Intervention", "Effort", "Impact", "Timeline")]
            for ai in dept.ai_interventions[:6]:
                rows.append((ai.get("name", ""), ai.get("effort", "").upper(),
                             ai.get("impact", "").upper(), ai.get("timeline", "")))
            _add_table(slide, rows, Inches(0.4), Inches(4.2), Inches(12.5),
                       col_widths=[Inches(4.5), Inches(2), Inches(2), Inches(4)],
                       row_font_size=9)

    # ------------------------------------------------------------------
    # SLIDE 9: Role Transformation Roadmap
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _slide_title_bar(slide, "Role Transformation Roadmap",
                     f"{len(blueprint.role_transformations)} roles identified for transformation")

    if blueprint.role_transformations:
        rows = [("Current Role", "Future Role", "Department", "Timeline", "Risk")]
        for rt in blueprint.role_transformations:
            rows.append((rt.current_role, rt.future_role, rt.department,
                         f"{rt.timeline_months} months", rt.risk_level.upper()))
        _add_table(slide, rows, Inches(0.4), Inches(1.3), Inches(12.5),
                   col_widths=[Inches(2.8), Inches(3.2), Inches(2.5), Inches(2), Inches(2)],
                   row_font_size=10)

    # ------------------------------------------------------------------
    # SLIDE 10: Skills Gap Summary
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _slide_title_bar(slide, "Skills Gap Analysis",
                     "Training requirements across all role transformations")

    if blueprint.skills_gap:
        sg = blueprint.skills_gap
        _add_metric_box(slide, Inches(0.4), Inches(1.3), Inches(3), Inches(0.9),
                        "Total Training Hours", f"{sg.total_training_hours:,.0f}")
        _add_metric_box(slide, Inches(3.8), Inches(1.3), Inches(3), Inches(0.9),
                        "Total Training Cost", f"${sg.total_training_cost:,.0f}")
        _add_metric_box(slide, Inches(7.2), Inches(1.3), Inches(3), Inches(0.9),
                        "Critical Gaps", str(len(sg.critical_gaps)), colour=_RED_RISK)
        _add_metric_box(slide, Inches(10.6), Inches(1.3), Inches(2.3), Inches(0.9),
                        "Timeline", f"{sg.training_timeline_months}mo")

        if sg.critical_gaps:
            rows = [("Role", "Skill Gap", "Urgency", "Department")]
            for gap in sg.critical_gaps[:8]:
                rows.append((gap.get("role", ""), gap.get("skill", ""),
                             gap.get("urgency", "").upper(), gap.get("department", "")))
            _add_table(slide, rows, Inches(0.4), Inches(2.6), Inches(12.5),
                       col_widths=[Inches(3), Inches(4), Inches(2), Inches(3.5)],
                       row_font_size=10)

    # ------------------------------------------------------------------
    # SLIDE 11: Change Management Plan
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _slide_title_bar(slide, "Change Management Plan",
                     "Phased approach to AI adoption")

    if blueprint.change_management:
        cm = blueprint.change_management
        y = Inches(1.3)
        for phase in cm.phases[:4]:
            lines = [f"- {obj[:65]}" for obj in phase.get("objectives", [])[:3]]
            _add_card(slide, Inches(0.4), y, Inches(6), Inches(1.3),
                      f"{phase.get('name', '')} ({phase.get('duration_months', '')} months)",
                      lines, accent_colour=_GREEN_ACC)
            y += Inches(1.45)

        # Resistance risks
        y = Inches(1.3)
        for risk in cm.resistance_mitigation[:4]:
            _add_card(slide, Inches(6.8), y, Inches(6), Inches(1.3),
                      f"RISK: {risk.get('risk', '')}",
                      [f"Likelihood: {risk.get('likelihood', '')}", f"Mitigation: {risk.get('mitigation', '')[:80]}"],
                      accent_colour=_ORANGE)
            y += Inches(1.45)

    # ------------------------------------------------------------------
    # SLIDE 12: Implementation Timeline
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _slide_title_bar(slide, "12-Month Implementation Timeline",
                     "Phased rollout with investment breakdown")

    if blueprint.implementation_phases:
        rows = [("Phase", "Months", "Focus Areas", "Investment")]
        for phase in blueprint.implementation_phases:
            focus = ", ".join(phase.get("focus_areas", [])[:3])
            rows.append((phase.get("phase", ""), phase.get("months", ""),
                         focus[:50], f"${phase.get('investment', 0):,.0f}"))
        _add_table(slide, rows, Inches(0.4), Inches(1.5), Inches(12.5),
                   col_widths=[Inches(2.5), Inches(1.5), Inches(6), Inches(2.5)],
                   row_font_size=11)

    # ------------------------------------------------------------------
    # SLIDE 13: Investment & ROI Summary
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _slide_title_bar(slide, "Investment & ROI Summary",
                     "Financial case for AI transformation")

    _add_metric_box(slide, Inches(0.4), Inches(1.5), Inches(3), Inches(1.1),
                    "Total Investment", f"${blueprint.total_implementation_cost:,.0f}")
    _add_metric_box(slide, Inches(3.8), Inches(1.5), Inches(3), Inches(1.1),
                    "Annual Savings", f"${blueprint.total_annual_savings:,.0f}", colour=_GREEN_ACC)
    _add_metric_box(slide, Inches(7.2), Inches(1.5), Inches(3), Inches(1.1),
                    "ROI Multiplier", f"{blueprint.roi_multiplier}x", colour=_GREEN_ACC)
    _add_metric_box(slide, Inches(10.6), Inches(1.5), Inches(2.3), Inches(1.1),
                    "Payback", f"{blueprint.payback_months}mo", colour=_GREEN_ACC)

    # Workflow savings breakdown
    if blueprint.workflow_redesigns:
        rows = [("Workflow", "Hours Saved/Year", "Time Savings %", "Implementation Cost")]
        for wf in blueprint.workflow_redesigns:
            rows.append((wf.workflow_name, f"{wf.annual_hours_saved:,.0f}",
                         f"{wf.time_savings_percent}%", f"${wf.implementation_cost:,.0f}"))
        _add_table(slide, rows, Inches(0.4), Inches(3.0), Inches(12.5),
                   col_widths=[Inches(3.5), Inches(3), Inches(3), Inches(3)],
                   row_font_size=10)

    # Save
    prs.save(str(output_path))
    return output_path
